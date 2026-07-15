#!/usr/bin/env python3
"""
Generate a Technohub-branded PPTX file from product data.
Usage: python3 generate_pptx.py <input_json> <output_pptx>
"""

import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import io

# EMU (English Metric Units) to inches conversion
EMU_PER_INCH = 914400

def emu_to_inches(emu):
    return emu / EMU_PER_INCH

def load_template_specs():
    """Load template specifications."""
    specs = {
        "slide": {
            "width_emu": 10691813,
            "height_emu": 7559675,
        },
        "card": {
            "width_emu": 2554041,
            "height_emu": 1459469,
            "columns_per_row": 4,
            "rows_per_slide": 5,
            "margin_left_emu": 266700,
            "margin_top_emu": 190914,
        },
        "logo": {
            "width_emu": 2414041,
            "height_emu": 460000,
            "position_left_emu": 70000,
            "position_top_emu": 40000,
        },
        "line": {
            "width_emu": 2554041,
            "height_emu": 60000,
            "position_left_emu": 0,
            "position_top_emu": 520000,
            "color_hex": "4472C4",
        },
        "text_elements": {
            "lao_name": {
                "font_size_pt": 11,
                "font_bold": True,
                "color_hex": "000000",
                "position_left_emu": 70000,
                "position_top_emu": 650000,
                "width_emu": 2414041,
                "height_emu": 380000, # Combined height for Lao + English
            },
            "english_name": {
                "font_size_pt": 8.5,
                "font_bold": False,
                "color_hex": "6499DE",
            },
            "price": {
                "font_size_pt": 18,
                "font_bold": True,
                "color_hex": "FF0000",
                "position_left_emu": 1000000,
                "position_top_emu": 1100000,
                "width_emu": 1484041,
                "height_emu": 280000,
            },
            "barcode": {
                "font_size_pt": 9,
                "font_bold": False,
                "color_hex": "000000",
                "position_left_emu": 70000,
                "position_top_emu": 1280000,
                "width_emu": 1000000,
                "height_emu": 150000,
            },
        },
        "border": {
            "color_hex": "000000",
            "width_pt": 2.0,
        },
    }
    return specs

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def add_card_border(slide, left, top, width, height, border_color_hex):
    """Add a border rectangle around the card."""
    border_color = RGBColor(*hex_to_rgb(border_color_hex))
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left,
        top,
        width,
        height,
    )
    shape.fill.background()
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)

def add_card_to_slide(slide, product, specs, card_index, logo_path):
    """Add a single product card to the slide."""
    card_spec = specs["card"]
    
    # Calculate position based on card index
    col = card_index % card_spec["columns_per_row"]
    row = (card_index // card_spec["columns_per_row"]) % card_spec["rows_per_slide"]
    
    left = card_spec["margin_left_emu"] + col * card_spec["width_emu"]
    top = card_spec["margin_top_emu"] + row * card_spec["height_emu"]
    
    # Add card border
    add_card_border(
        slide,
        left,
        top,
        card_spec["width_emu"],
        card_spec["height_emu"],
        specs["border"]["color_hex"],
    )
    
    # Add logo at the top of the card
    logo_spec = specs["logo"]
    logo_left = left + logo_spec["position_left_emu"]
    logo_top = top + logo_spec["position_top_emu"]
    
    try:
        slide.shapes.add_picture(
            logo_path,
            logo_left,
            logo_top,
            width=logo_spec["width_emu"],
            height=logo_spec["height_emu"],
        )
    except Exception as e:
        print(f"Warning: Could not add logo: {e}", file=sys.stderr)
        
    # Add horizontal separator line
    line_spec = specs["line"]
    line_left = left + line_spec["position_left_emu"]
    line_top = top + line_spec["position_top_emu"]
    try:
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            line_left,
            line_top,
            line_spec["width_emu"],
            line_spec["height_emu"],
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(line_spec["color_hex"]))
        line_shape.line.color.rgb = RGBColor(*hex_to_rgb(line_spec["color_hex"]))
        line_shape.line.width = Pt(1)
    except Exception as e:
        print(f"Warning: Could not add line: {e}", file=sys.stderr)
    
    # Add Lao and English names in a single textbox to auto-flow properly
    lao_spec = specs["text_elements"]["lao_name"]
    english_spec = specs["text_elements"]["english_name"]
    name_left = left + lao_spec["position_left_emu"]
    name_top = top + lao_spec["position_top_emu"]
    
    # The height is now taken directly from lao_spec since it defines the combined bounding box
    name_height = lao_spec["height_emu"]
    
    name_box = slide.shapes.add_textbox(
        name_left, name_top, lao_spec["width_emu"], name_height
    )
    name_frame = name_box.text_frame
    name_frame.word_wrap = True
    name_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    name_frame.vertical_anchor = MSO_ANCHOR.TOP
    name_frame.margin_left = Inches(0)
    name_frame.margin_right = Inches(0)
    name_frame.margin_top = Inches(0)
    name_frame.margin_bottom = Inches(0)
    
    # Paragraph 1: Lao Name
    p_lao = name_frame.paragraphs[0]
    p_lao.text = product["laoName"]
    p_lao.font.name = "Noto Sans Lao"
    p_lao.font.size = Pt(lao_spec["font_size_pt"])
    if lao_spec.get("font_bold"):
        p_lao.font.bold = True
    p_lao.font.color.rgb = RGBColor(*hex_to_rgb(lao_spec["color_hex"]))
    
    # Paragraph 2: English Name
    if product["englishName"]:
        p_eng = name_frame.add_paragraph()
        p_eng.text = product["englishName"]
        p_eng.font.name = "Noto Sans"
        p_eng.font.size = Pt(english_spec["font_size_pt"])
        if english_spec.get("font_bold"):
            p_eng.font.bold = True
        p_eng.font.color.rgb = RGBColor(*hex_to_rgb(english_spec["color_hex"]))
    
    # Add Price
    price_spec = specs["text_elements"]["price"]
    price_left = left + price_spec["position_left_emu"]
    price_top = top + price_spec["position_top_emu"]
    price_box = slide.shapes.add_textbox(
        price_left, price_top, price_spec["width_emu"], price_spec["height_emu"]
    )
    price_frame = price_box.text_frame
    price_frame.word_wrap = False
    price_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    price_frame.vertical_anchor = MSO_ANCHOR.TOP
    price_frame.margin_left = Inches(0)
    price_frame.margin_right = Inches(0)
    price_frame.margin_top = Inches(0)
    price_frame.margin_bottom = Inches(0)
    price_p = price_frame.paragraphs[0]
    price_p.text = product["price"]
    price_p.font.name = "Noto Sans"
    price_p.font.size = Pt(price_spec["font_size_pt"])
    price_p.font.bold = price_spec["font_bold"]
    price_p.font.color.rgb = RGBColor(*hex_to_rgb(price_spec["color_hex"]))
    price_p.alignment = PP_ALIGN.RIGHT
    
    # Add Barcode
    barcode_spec = specs["text_elements"]["barcode"]
    barcode_left = left + barcode_spec["position_left_emu"]
    barcode_top = top + barcode_spec["position_top_emu"]
    barcode_box = slide.shapes.add_textbox(
        barcode_left, barcode_top, barcode_spec["width_emu"], barcode_spec["height_emu"]
    )
    barcode_frame = barcode_box.text_frame
    barcode_frame.word_wrap = True
    barcode_frame.vertical_anchor = MSO_ANCHOR.TOP
    barcode_frame.margin_left = Inches(0)
    barcode_frame.margin_right = Inches(0)
    barcode_frame.margin_top = Inches(0)
    barcode_frame.margin_bottom = Inches(0)
    barcode_p = barcode_frame.paragraphs[0]
    barcode_p.text = product["barcode"]
    barcode_p.font.name = "Noto Sans"
    barcode_p.font.size = Pt(barcode_spec["font_size_pt"])
    barcode_p.font.color.rgb = RGBColor(*hex_to_rgb(barcode_spec["color_hex"]))

def generate_pptx(products, output_path, logo_path):
    """Generate PPTX from product data."""
    specs = load_template_specs()
    
    # Create presentation with custom slide size
    prs = Presentation()
    prs.slide_width = specs["slide"]["width_emu"]
    prs.slide_height = specs["slide"]["height_emu"]
    
    # Add blank slides and populate with cards
    cards_per_slide = specs["card"]["columns_per_row"] * specs["card"]["rows_per_slide"]
    num_slides = (len(products) + cards_per_slide - 1) // cards_per_slide
    
    for slide_idx in range(num_slides):
        # Add blank slide
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add white background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Add cards to this slide
        start_idx = slide_idx * cards_per_slide
        end_idx = min(start_idx + cards_per_slide, len(products))
        
        for card_idx in range(end_idx - start_idx):
            product = products[start_idx + card_idx]
            add_card_to_slide(slide, product, specs, card_idx, logo_path)
    
    # Save presentation
    prs.save(output_path)

def main():
    if len(sys.argv) != 3:
        print("Usage: python3 generate_pptx.py <input_json> <output_pptx>", file=sys.stderr)
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_file = sys.argv[2]
    
    # Determine logo path
    script_dir = Path(__file__).parent
    logo_path = script_dir / "assets" / "techno-hub-logo.jpg"
    
    if not logo_path.exists():
        print(f"Error: Logo file not found at {logo_path}", file=sys.stderr)
        sys.exit(1)
    
    try:
        with open(input_file, "r", encoding="utf-8") as f:
            products = json.load(f)
        
        generate_pptx(products, output_file, str(logo_path))
        print(f"PPTX generated successfully: {output_file}")
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()
